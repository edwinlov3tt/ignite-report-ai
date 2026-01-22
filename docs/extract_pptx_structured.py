#!/usr/bin/env python3
"""
Extract PowerPoint content into structured JSON format for AI analysis.
Preserves slide structure, tables, and text hierarchy.
"""

import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_table(table):
    """Extract table data as list of dicts with headers as keys."""
    rows = []
    headers = []
    
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        if i == 0:
            headers = cells
        else:
            if headers:
                row_dict = {headers[j]: cells[j] for j in range(min(len(headers), len(cells)))}
            else:
                row_dict = {f"col_{j}": cells[j] for j in range(len(cells))}
            rows.append(row_dict)
    
    return {"headers": headers, "rows": rows}

def extract_text_frame(text_frame):
    """Extract text from a text frame with paragraph structure."""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append({
                "text": text,
                "level": para.level
            })
    return paragraphs

def extract_presentation(filepath):
    """Extract all content from a PowerPoint presentation."""
    prs = Presentation(filepath)
    
    presentation_data = {
        "metadata": {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width.inches if prs.slide_width else None,
            "slide_height": prs.slide_height.inches if prs.slide_height else None
        },
        "slides": []
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "title": None,
            "content": [],
            "tables": [],
            "notes": None
        }
        
        # Extract title if present
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()
        
        # Extract all shapes
        for shape in slide.shapes:
            # Skip title shape (already captured)
            if shape == slide.shapes.title:
                continue
            
            # Handle tables
            if shape.has_table:
                table_data = extract_table(shape.table)
                if table_data["rows"]:  # Only add non-empty tables
                    slide_data["tables"].append(table_data)
            
            # Handle text frames
            elif shape.has_text_frame:
                paragraphs = extract_text_frame(shape.text_frame)
                for para in paragraphs:
                    # Skip image placeholders and backgrounds
                    if not para["text"].startswith(("/src/", "preencoded", "Image")):
                        slide_data["content"].append(para)
        
        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text
        
        # Clean up empty fields
        if not slide_data["content"]:
            del slide_data["content"]
        if not slide_data["tables"]:
            del slide_data["tables"]
        if not slide_data["notes"]:
            del slide_data["notes"]
        
        presentation_data["slides"].append(slide_data)
    
    return presentation_data

def main():
    filepath = "/mnt/user-data/uploads/1768506094356_lumina_charts.pptx"
    
    print("Extracting PowerPoint content...")
    data = extract_presentation(filepath)
    
    # Save as JSON
    output_path = "/home/claude/pptx_extraction_structured.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    
    print(f"Extracted {data['metadata']['slide_count']} slides")
    print(f"Saved to: {output_path}")
    
    # Also create a simplified markdown version for quick reading
    md_output = "/home/claude/pptx_extraction_clean.md"
    with open(md_output, 'w', encoding='utf-8') as f:
        f.write(f"# PowerPoint Extraction\n\n")
        f.write(f"**Total Slides:** {data['metadata']['slide_count']}\n\n")
        f.write("---\n\n")
        
        for slide in data["slides"]:
            f.write(f"## Slide {slide['slide_number']}")
            if slide.get("title"):
                f.write(f": {slide['title']}")
            f.write("\n\n")
            
            # Write content
            if slide.get("content"):
                for item in slide["content"]:
                    indent = "  " * item.get("level", 0)
                    f.write(f"{indent}- {item['text']}\n")
                f.write("\n")
            
            # Write tables
            if slide.get("tables"):
                for i, table in enumerate(slide["tables"]):
                    if table["headers"]:
                        f.write("| " + " | ".join(table["headers"]) + " |\n")
                        f.write("| " + " | ".join(["---"] * len(table["headers"])) + " |\n")
                        for row in table["rows"]:
                            values = [str(row.get(h, "")) for h in table["headers"]]
                            f.write("| " + " | ".join(values) + " |\n")
                        f.write("\n")
            
            f.write("---\n\n")
    
    print(f"Clean markdown saved to: {md_output}")

if __name__ == "__main__":
    main()
