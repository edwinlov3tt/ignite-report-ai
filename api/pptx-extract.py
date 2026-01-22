"""
PowerPoint Table Extraction API
Extracts tables from PPTX files and returns structured JSON.
"""

import json
import io
import cgi
from http.server import BaseHTTPRequestHandler
from pptx import Presentation

# Constants
MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50MB
PPTX_MAGIC_BYTES = b'PK'  # PPTX files are ZIP archives starting with PK


def extract_table(table):
    """Extract table data as list of dicts with headers as keys.

    Args:
        table: A python-pptx Table object

    Returns:
        dict with 'headers' (list) and 'rows' (list of dicts)
    """
    rows = []
    headers = []

    for i, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        if i == 0:
            # First row is headers
            headers = cells
        else:
            if headers:
                row_dict = {headers[j]: cells[j] for j in range(min(len(headers), len(cells)))}
            else:
                # Fallback for tables without headers
                row_dict = {f"Column {j + 1}": cells[j] for j in range(len(cells))}
            rows.append(row_dict)

    return {"headers": headers, "rows": rows}


def extract_presentation(file_bytes: bytes) -> dict:
    """Extract all tables from a PowerPoint presentation.

    Args:
        file_bytes: Raw bytes of the PPTX file

    Returns:
        dict with 'metadata' and 'tables' keys
    """
    prs = Presentation(io.BytesIO(file_bytes))

    tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # Get slide title
        slide_title = None
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()

        # Track table index within slide
        table_index = 0

        # Extract all tables from slide
        for shape in slide.shapes:
            if shape.has_table:
                table_data = extract_table(shape.table)

                # Only include non-empty tables
                if table_data["rows"]:
                    tables.append({
                        "slideNumber": slide_num,
                        "slideTitle": slide_title,
                        "tableIndex": table_index,
                        "headers": table_data["headers"],
                        "rows": table_data["rows"],
                        "rowCount": len(table_data["rows"])
                    })
                    table_index += 1

    return {
        "metadata": {
            "slideCount": len(prs.slides)
        },
        "tables": tables
    }


def validate_pptx_file(file_bytes: bytes, filename: str) -> str | None:
    """Validate that the uploaded file is a valid PPTX.

    Args:
        file_bytes: Raw bytes of the file
        filename: Original filename

    Returns:
        Error message if invalid, None if valid
    """
    # Check file extension
    if not filename.lower().endswith('.pptx'):
        return "Please upload a .pptx file"

    # Check file size
    if len(file_bytes) > MAX_FILE_SIZE_BYTES:
        return f"File size must be less than {MAX_FILE_SIZE_BYTES // (1024 * 1024)}MB"

    # Check magic bytes (PPTX files are ZIP archives)
    if len(file_bytes) < 2 or file_bytes[:2] != PPTX_MAGIC_BYTES:
        return "Invalid PowerPoint file format"

    return None


class handler(BaseHTTPRequestHandler):
    """Vercel serverless function handler for PPTX extraction."""

    def do_POST(self):
        try:
            # Parse content type
            content_type = self.headers.get('Content-Type', '')

            if 'multipart/form-data' not in content_type:
                self.send_error_response(400, "Expected multipart/form-data")
                return

            # Get content length
            content_length = int(self.headers.get('Content-Length', 0))
            if content_length == 0:
                self.send_error_response(400, "No file uploaded")
                return

            # Parse multipart form data using cgi module
            environ = {
                'REQUEST_METHOD': 'POST',
                'CONTENT_TYPE': content_type,
                'CONTENT_LENGTH': str(content_length),
            }

            form = cgi.FieldStorage(
                fp=self.rfile,
                headers=self.headers,
                environ=environ
            )

            # Get the file from form data
            file_item = form.get('file')
            if file_item is None or not hasattr(file_item, 'file'):
                self.send_error_response(400, "No file found in request")
                return

            # Read file content
            file_bytes = file_item.file.read()
            filename = file_item.filename or "unknown.pptx"

            # Validate the file
            validation_error = validate_pptx_file(file_bytes, filename)
            if validation_error:
                self.send_error_response(400, validation_error)
                return

            # Extract presentation
            result = extract_presentation(file_bytes)
            result["success"] = True
            result["metadata"]["filename"] = filename

            if len(result["tables"]) == 0:
                result["warning"] = "No data tables found in this PowerPoint"

            # Send success response
            self.send_response(200)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            self.wfile.write(json.dumps(result).encode())

        except Exception as e:
            self.send_error_response(500, f"Failed to extract data: {str(e)}")

    def do_OPTIONS(self):
        """Handle CORS preflight requests."""
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def send_error_response(self, status_code: int, message: str):
        """Send a JSON error response.

        Args:
            status_code: HTTP status code
            message: Error message to include
        """
        self.send_response(status_code)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        response = {"success": False, "error": message}
        self.wfile.write(json.dumps(response).encode())
